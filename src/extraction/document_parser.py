"""
Document parser module for extracting text from PDFs, PPTs, and Word documents.
"""

import logging
from pathlib import Path
from typing import Dict, List, Tuple
import PyPDF2
from pptx import Presentation
from docx import Document

logger = logging.getLogger(__name__)


class DocumentParser:
    """Extract text from various document formats."""

    @staticmethod
    def extract_from_pdf(file_path: str) -> Tuple[str, List[Dict]]:
        """
        Extract text from PDF file with page information.
        
        Returns:
            Tuple of (full_text, pages_info)
            pages_info contains: [{"page_num": int, "text": str}, ...]
        """
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                full_text = ""
                pages_info = []

                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    full_text += text + "\n"
                    pages_info.append({
                        "page_num": page_num,
                        "text": text
                    })

            logger.info(f"Successfully extracted text from PDF: {file_path}")
            return full_text, pages_info
        except Exception as e:
            logger.error(f"Error extracting from PDF {file_path}: {str(e)}")
            raise

    @staticmethod
    def extract_from_ppt(file_path: str) -> Tuple[str, List[Dict]]:
        """
        Extract text from PowerPoint file with slide information.
        
        Returns:
            Tuple of (full_text, slides_info)
            slides_info contains: [{"slide_num": int, "text": str}, ...]
        """
        try:
            presentation = Presentation(file_path)
            full_text = ""
            slides_info = []

            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"

                full_text += slide_text + "\n"
                slides_info.append({
                    "slide_num": slide_num,
                    "text": slide_text
                })

            logger.info(f"Successfully extracted text from PPT: {file_path}")
            return full_text, slides_info
        except Exception as e:
            logger.error(f"Error extracting from PPT {file_path}: {str(e)}")
            raise

    @staticmethod
    def extract_from_word(file_path: str) -> Tuple[str, List[Dict]]:
        """
        Extract text from Word document.
        
        Returns:
            Tuple of (full_text, paragraphs_info)
            paragraphs_info contains: [{"para_num": int, "text": str}, ...]
        """
        try:
            doc = Document(file_path)
            full_text = ""
            paragraphs_info = []

            for para_num, para in enumerate(doc.paragraphs, 1):
                if para.text.strip():  # Only include non-empty paragraphs
                    full_text += para.text + "\n"
                    paragraphs_info.append({
                        "para_num": para_num,
                        "text": para.text
                    })

            logger.info(f"Successfully extracted text from Word: {file_path}")
            return full_text, paragraphs_info
        except Exception as e:
            logger.error(f"Error extracting from Word {file_path}: {str(e)}")
            raise

    @staticmethod
    def extract_from_file(file_path: str) -> Tuple[str, List[Dict], str]:
        """
        Auto-detect file type and extract text accordingly.
        
        Returns:
            Tuple of (full_text, content_info, file_type)
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        if extension == ".pdf":
            text, info = DocumentParser.extract_from_pdf(str(file_path))
            return text, info, "pdf"
        elif extension in [".ppt", ".pptx"]:
            text, info = DocumentParser.extract_from_ppt(str(file_path))
            return text, info, "ppt"
        elif extension in [".doc", ".docx"]:
            text, info = DocumentParser.extract_from_word(str(file_path))
            return text, info, "word"
        else:
            raise ValueError(f"Unsupported file format: {extension}")


def extract_table_of_contents(text: str) -> List[Tuple[str, int]]:
    """
    Extract Table of Contents from document text.
    
    Returns:
        List of (section_title, approximate_line_number)
    """
    toc_items = []
    lines = text.split('\n')
    
    in_toc = False
    for line_num, line in enumerate(lines):
        # Detect start of TOC
        if "table of contents" in line.lower() or "contents" in line.lower():
            in_toc = True
            continue
        
        # Simple detection: lines with numbers followed by text (common TOC format)
        if in_toc:
            stripped = line.strip()
            # Stop at certain patterns that indicate end of TOC
            if not stripped or (len(stripped) < 5 and not any(c.isalpha() for c in stripped)):
                if len(toc_items) > 3:  # Consider it end of TOC if we found several items
                    break
            
            # Extract section titles (usually start with number or bullet)
            if stripped and (stripped[0].isdigit() or stripped[0] in ['•', '-', '*']):
                toc_items.append((stripped, line_num))
    
    logger.info(f"Extracted {len(toc_items)} items from Table of Contents")
    return toc_items
