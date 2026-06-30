import re
import logging
import json
import uuid
from pathlib import Path
from typing import Dict, List, Tuple, Optional
import faiss
import numpy as np
# NOTE: PyPDF2 / python-pptx / python-docx are imported lazily inside the
# _extract_* helpers below. They are only needed to parse source documents when
# (re)building the proposal index, never to serve queries — keeping them out of
# the import path keeps the runtime (serving) memory footprint small.

from config import (
    RAW_PROPOSALS_DIR,
    PROPOSAL_FAISS_INDEX_PATH,
    PROPOSAL_METADATA_PATH,
    PROPOSAL_VECTOR_BACKEND,
    QDRANT_URL,
    QDRANT_API_KEY,
    QDRANT_PROPOSAL_COLLECTION,
    get_embeddings,
)

# Stable namespace so a chunk (source_file + chunk_index) always maps to the same
# Qdrant point ID → re-embedding a doc overwrites its chunks in place.
_QDRANT_NS = uuid.UUID("6f9619ff-8b86-d011-b42d-00cf4fc964ff")

logger = logging.getLogger(__name__)

# ── Chunking knobs ──────────────────────────────────────────────────────────────
# A section whose body is at/under MAX_WORDS stays one chunk; larger leaf bodies
# are split at paragraph boundaries. Consecutive sibling sections each under
# MERGE_WORDS are merged into one chunk (avoids a swarm of tiny vectors).
MAX_WORDS = 350
MERGE_WORDS = 80
MIN_CHUNK_CHARS = 60

# ── Boilerplate paragraph / line detection ──────────────────────────────────────
_BOILERPLATE_RE = re.compile(
    r'\.{5,}'                               # TOC dots
    r'|Page \| \d+'                         # "Page | 12"
    r'|Copyright ©'                         # footer
    r'|Axestrack Software Solutions Pvt'    # address block
    r'|Gopalpura bypass'
    r'|Office No\.'
    r'|Submitted on\s*[–\-]'               # cover page date
    r'|Table of Contents'
    r'|^\s*\d+\s*$',                        # lone page numbers
    re.IGNORECASE | re.MULTILINE
)

# ── Heading detection ───────────────────────────────────────────────────────────
# Numbered headings like "1", "1.2", "2.2.1 Smart Traffic ...". The dotted number
# gives both the path and the depth (number of components). Title must look like a
# heading: starts with a letter, not too long, and not ending in sentence period.
_NUM_HEADING_RE = re.compile(r'^\s*(\d+(?:\.\d+)*)\.?\s+([A-Za-z(].{0,90})$')


def _numbered_heading(line: str) -> Optional[Tuple[str, str, int]]:
    """Return (path, title, level) if the line is a numbered heading, else None."""
    m = _NUM_HEADING_RE.match(line)
    if not m:
        return None
    path, title = m.group(1), m.group(2).strip()
    # Reject lines that are really sentences/bullets that happen to start with a number.
    if title.endswith('.') and len(title.split()) > 12:
        return None
    if title.endswith((',', ';', ':')) and len(title.split()) > 12:
        return None
    level = path.count('.') + 1
    return path, title, level


def _good_top_title(title: str) -> bool:
    """A plausible top-level section title: no digits, starts uppercase, and is
    either multi-word or a single long non-acronym word ('Assumptions'). Rejects
    table cells like 'CRM', 'MP', 'GB + 64 GB' that start with a stray integer."""
    if any(ch.isdigit() for ch in title):
        return False
    if not title[:1].isupper():
        return False
    alpha = re.findall(r"[A-Za-z]+", title)
    if not alpha:
        return False
    if len(alpha) >= 2:
        return True
    w = alpha[0]
    return len(w) >= 5 and not w.isupper()


def _validate_headings(elements: List[Dict]) -> None:
    """Reject numbered-heading false positives from numeric tables (rows like
    '4 GB + 64 GB', '14 Hypercare & Support'). Walks elements in document order:
    a decimal heading is valid only if its parent heading was already seen; a bare
    integer N is valid only if it continues the section sequence (max_top < N <=
    max_top+2, tolerating one image-only/missing heading) and has a real title.
    Invalid headings are demoted back to body text so their content isn't lost.
    Only PDF (numbered) headings are checked — Word/PPT headings have path=None."""
    # A page with many bare-integer "headings" is a numeric table (Gantt, BOQ, spec
    # sheet), not a run of sections — its integer rows are demoted wholesale.
    DENSE = 4
    page_int_count: Dict[str, int] = {}
    for el in elements:
        if el.get("kind") == "heading" and el.get("path") and '.' not in el["path"]:
            pr = el.get("page_ref", "")
            page_int_count[pr] = page_int_count.get(pr, 0) + 1

    seen = set()
    max_top = 0
    for el in elements:
        if el.get("kind") != "heading" or not el.get("path"):
            continue
        path, title = el["path"], el["text"]
        if '.' in path:
            valid = path.rsplit('.', 1)[0] in seen
        else:
            n = int(path)
            dense = page_int_count.get(el.get("page_ref", ""), 0) >= DENSE
            valid = (not dense) and (max_top < n <= max_top + 2) and _good_top_title(title)
        if valid:
            seen.add(path)
            if '.' not in path:
                max_top = int(path)
        else:
            el["kind"] = "body"
            el["text"] = f"{path} {title}"


# ── Industry classification ─────────────────────────────────────────────────────
_CLIENT_INDUSTRY = [
    ('amns',           'steel'),
    ('arcelormittal',  'steel'),
    ('kattupalli',     'port'),
    ('ennore',         'port'),
    ('abg trading',    'port'),
    ('bajel',          'port'),
    ('ashok leyland',  'automotive'),
    ('ather',          'automotive'),
    ('apcotex',        'chemical'),
    ('atul',           'chemical'),
    ('amar raja',      'battery'),
    ('amazon',         'ecommerce'),
    ('asahi',          'glass'),
    ('ashai',          'glass'),
    ('adani',          'conglomerate'),
]

_CONTENT_INDUSTRY = {
    'steel':      ['steel plant', 'hazira', 'marshalling yard', 'rolling mill'],
    'port':       ['vessel', 'berth', 'barge', 'jetty', 'coal shipment', 'navis', 'stevedoring'],
    'automotive': ['electric vehicle', 'automobile', 'oem'],
    'chemical':   ['chemical', 'polymer', 'latex', 'hazmat', 'hazardous material'],
    'ecommerce':  ['last-mile', 'last mile', 'delivery hub', 'e-commerce'],
    'battery':    ['battery', 'energy storage'],
    'glass':      ['glass'],
    'mining':     ['mine', 'mining', 'mineral'],
}


def _detect_industry(text: str, client_name: str) -> str:
    cl = client_name.lower()
    for name, industry in _CLIENT_INDUSTRY:
        if name in cl:
            return industry
    body = text[:2000].lower()
    for industry, kws in _CONTENT_INDUSTRY.items():
        for kw in kws:
            if re.search(r'\b' + re.escape(kw) + r'\b', body):
                return industry
    return 'logistics'


def _is_boilerplate(para: str) -> bool:
    if _BOILERPLATE_RE.search(para):
        return True
    if 'axestrack has been recognized' in para.lower() and 'gartner' in para.lower():
        return True
    return False


def _clean_text(text: str) -> str:
    text = re.sub(r'[ \t]{3,}', '  ', text)
    text = re.sub(r'-\s*\n\s*', '', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


# ── Section model ───────────────────────────────────────────────────────────────
# A Section is one heading plus the body text directly under it (before the next
# heading). `path` is the dotted number ("2.2.1"), `parent` is the path with the
# last component dropped ("2.2", or None at the top level). `breadcrumb` is the
# chain of ancestor titles. Sections are produced in document order.
class Section:
    __slots__ = ("path", "title", "level", "parent", "breadcrumb", "body", "page_ref")

    def __init__(self, path, title, level, parent, breadcrumb, page_ref):
        self.path = path
        self.title = title
        self.level = level
        self.parent = parent
        self.breadcrumb = breadcrumb
        self.body: List[str] = []
        self.page_ref = page_ref


def _parent_of(path: str) -> Optional[str]:
    return path.rsplit('.', 1)[0] if '.' in path else None


# ── Format extractors → ordered list of "elements" ──────────────────────────────
# Each element is a dict: {kind: 'heading'|'body', text, page_ref, level?}.
# A common assembler turns elements into Sections, so the chunking logic is shared
# across PDF / PPT / Word.

def _extract_pdf_elements(file_path: Path) -> List[Dict]:
    try:
        from PyPDF2 import PdfReader            # production dependency
    except ImportError:
        from pypdf import PdfReader             # successor package, identical API
    elements: List[Dict] = []
    try:
        with open(file_path, 'rb') as f:
            reader = PdfReader(f)
            for i, page in enumerate(reader.pages, 1):
                raw = page.extract_text() or ''
                if not raw.strip():
                    continue
                page_ref = f"p.{i}"
                # Split into lines; numbered lines are headings, the rest is body.
                buf: List[str] = []
                for line in _clean_text(raw).split('\n'):
                    line = line.strip()
                    if not line or _is_boilerplate(line):
                        continue
                    h = _numbered_heading(line)
                    if h:
                        if buf:
                            elements.append({"kind": "body", "text": "\n".join(buf), "page_ref": page_ref})
                            buf = []
                        path, title, level = h
                        elements.append({"kind": "heading", "path": path, "text": title,
                                         "level": level, "page_ref": page_ref})
                    else:
                        buf.append(line)
                if buf:
                    elements.append({"kind": "body", "text": "\n".join(buf), "page_ref": page_ref})
    except Exception as e:
        logger.error(f"PDF extraction failed for {file_path.name}: {e}")
    return elements


def _extract_pptx_elements(file_path: Path) -> List[Dict]:
    """PPT has no nested hierarchy: each slide is one level-1 section whose title is
    the slide title (or 'Slide N') and whose body is the remaining shape text."""
    from pptx import Presentation
    elements: List[Dict] = []
    try:
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides, 1):
            title_text = ""
            try:
                if slide.shapes.title and slide.shapes.title.text.strip():
                    title_text = slide.shapes.title.text.strip()
            except Exception:
                title_text = ""
            # Body includes ALL shape text (title too) so title-only slides still
            # produce content; the title also doubles as the heading label.
            body_parts = []
            for shape in slide.shapes:
                if not hasattr(shape, 'text'):
                    continue
                t = shape.text.strip()
                if t:
                    body_parts.append(t)
            body = _clean_text("\n".join(body_parts))
            if not body:
                continue
            elements.append({"kind": "heading", "path": str(i),
                             "text": title_text or f"Slide {i}", "level": 1,
                             "page_ref": f"slide {i}"})
            elements.append({"kind": "body", "text": body, "page_ref": f"slide {i}"})
    except Exception as e:
        logger.error(f"PPT extraction failed for {file_path.name}: {e}")
    return elements


def _linearize_docx_table(table) -> str:
    """Flatten a Word table to 'col: val | col: val' rows so counts/specs stay
    queryable instead of becoming a wall of cells."""
    rows = []
    cells = [[c.text.strip() for c in row.cells] for row in table.rows]
    if not cells:
        return ""
    header = cells[0]
    looks_headed = all(h and len(h) < 40 for h in header) and len(header) > 1
    for r in cells[1:] if looks_headed else cells:
        if not any(r):
            continue
        if looks_headed:
            rows.append(" | ".join(f"{h}: {v}" for h, v in zip(header, r) if v))
        else:
            rows.append(" | ".join(v for v in r if v))
    return "\n".join(rows)


def _extract_docx_elements(file_path: Path) -> List[Dict]:
    """Word headings are explicit (style 'Heading N'), so we get a real leveled
    tree. Tables are linearized inline. Synthetic numbering is assigned by the
    assembler from the heading levels."""
    from docx import Document
    from docx.document import Document as _Doc
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    elements: List[Dict] = []
    try:
        doc = Document(file_path)
        parent_elm = doc.element.body
        for child in parent_elm.iterchildren():
            if child.tag.endswith('}p'):
                para = Paragraph(child, doc)
                text = para.text.strip()
                if not text or _is_boilerplate(text):
                    continue
                style = (para.style.name if para.style else '') or ''
                m = re.search(r'heading\s*(\d+)', style, re.IGNORECASE)
                if m:
                    level = int(m.group(1))
                    elements.append({"kind": "heading", "path": None, "text": text,
                                     "level": level, "page_ref": ""})
                elif style.lower() in ('title',):
                    elements.append({"kind": "heading", "path": None, "text": text,
                                     "level": 1, "page_ref": ""})
                else:
                    elements.append({"kind": "body", "text": text, "page_ref": ""})
            elif child.tag.endswith('}tbl'):
                table = Table(child, doc)
                flat = _linearize_docx_table(table)
                if flat:
                    elements.append({"kind": "body", "text": flat, "page_ref": ""})
    except Exception as e:
        logger.error(f"Word extraction failed for {file_path.name}: {e}")
    return elements


# ── Assemble elements → Sections → chunks ───────────────────────────────────────

def _assign_synthetic_paths(elements: List[Dict]) -> None:
    """For formats without numbered headings (Word, PPT), assign a dotted path per
    heading from running per-level counters so parent/sibling grouping still works.
    Numbered headings (PDF) keep their own path."""
    counters: Dict[int, int] = {}
    for el in elements:
        if el.get("kind") != "heading":
            continue
        if el.get("path"):
            continue  # numbered heading already has a real path
        level = el.get("level", 1)
        counters[level] = counters.get(level, 0) + 1
        for deeper in [k for k in counters if k > level]:
            counters.pop(deeper, None)
        el["path"] = ".".join(str(counters[k]) for k in sorted(counters) if k <= level)


def _build_sections(elements: List[Dict]) -> List[Section]:
    _assign_synthetic_paths(elements)
    sections: List[Section] = []
    title_by_path: Dict[str, str] = {}
    current: Optional[Section] = None
    # Body that appears before the first heading — or a whole document with no
    # detectable headings (scanned PDF, un-styled Word) — is collected under a
    # synthetic "Introduction" section so its content is chunked, not lost.
    for el in elements:
        if el["kind"] == "body" and current is None:
            current = Section("0", "Introduction", 1, None, "Introduction", el["page_ref"])
            title_by_path["0"] = "Introduction"
            sections.append(current)
            current.body.append(el["text"])
            continue
        if el["kind"] == "heading":
            path = el["path"]
            title = el["text"]
            title_by_path[path] = title
            parent = _parent_of(path)
            # breadcrumb from ancestor paths
            crumbs, p = [], path
            while p is not None:
                crumbs.append(title_by_path.get(p, ""))
                p = _parent_of(p)
            breadcrumb = " > ".join(c for c in reversed(crumbs) if c)
            current = Section(path, title, el["level"], parent, breadcrumb, el["page_ref"])
            sections.append(current)
        elif el["kind"] == "body" and current is not None:
            current.body.append(el["text"])
    return sections


def _split_body(text: str) -> List[str]:
    """Split an oversized body at paragraph boundaries into <= MAX_WORDS pieces."""
    paras = [p.strip() for p in re.split(r'\n{2,}', text) if p.strip()]
    if not paras:
        paras = [text]
    chunks, buf, wc = [], [], 0
    for p in paras:
        pw = len(p.split())
        if buf and wc + pw > MAX_WORDS:
            chunks.append("\n\n".join(buf))
            buf, wc = [], 0
        buf.append(p)
        wc += pw
    if buf:
        chunks.append("\n\n".join(buf))
    return chunks


def _has_descendants(section: Section, sections: List[Section]) -> bool:
    prefix = section.path + "."
    return any(s.path.startswith(prefix) for s in sections)


def _sections_to_chunks(sections: List[Section]) -> List[Dict]:
    """Turn Sections into chunk dicts, applying: drop empty heading-only parents,
    merge tiny consecutive siblings, split oversized leaf bodies."""
    raw = []
    for s in sections:
        body = _clean_text("\n\n".join(b for b in s.body if b.strip()))
        words = len(body.split())
        raw.append({
            "path": s.path, "title": s.title, "parent": s.parent,
            "breadcrumb": s.breadcrumb, "page_ref": s.page_ref,
            "body": body, "words": words,
            "has_children": _has_descendants(s, sections),
        })

    chunks: List[Dict] = []
    i = 0
    while i < len(raw):
        sec = raw[i]
        # Heading-only node that has children (e.g., "1.2") and no body: skip — its
        # children carry the content. If it has no children and no body, also skip.
        if not sec["body"]:
            i += 1
            continue

        # Merge a run of tiny consecutive siblings (same parent) into one chunk.
        if sec["words"] < MERGE_WORDS and not sec["has_children"]:
            group = [sec]
            j = i + 1
            while (j < len(raw) and raw[j]["body"] and not raw[j]["has_children"]
                   and raw[j]["words"] < MERGE_WORDS
                   and raw[j]["parent"] == sec["parent"]):
                group.append(raw[j])
                j += 1
            if len(group) > 1:
                merged_text = "\n\n".join(f"{g['title']}: {g['body']}" for g in group)
                chunks.append({
                    "path": sec["path"], "title": group[0]["title"] + f" (+{len(group)-1} more)",
                    "parent": sec["parent"], "breadcrumb": sec["breadcrumb"],
                    "page_ref": sec["page_ref"], "content": merged_text,
                })
                i = j
                continue

        # Oversized leaf → split; otherwise one chunk for the section.
        pieces = _split_body(sec["body"]) if sec["words"] > MAX_WORDS else [sec["body"]]
        for k, piece in enumerate(pieces):
            title = sec["title"] if len(pieces) == 1 else f"{sec['title']} (part {k+1})"
            chunks.append({
                "path": sec["path"], "title": title, "parent": sec["parent"],
                "breadcrumb": sec["breadcrumb"], "page_ref": sec["page_ref"],
                "content": piece,
            })
        i += 1
    return chunks


_EXTRACTORS = {
    '.pdf':  _extract_pdf_elements,
    '.ppt':  _extract_pptx_elements,
    '.pptx': _extract_pptx_elements,
    '.doc':  _extract_docx_elements,
    '.docx': _extract_docx_elements,
}

# The proposal index should hold Axestrack proposals/solution docs only — not the
# client's own RFQ requirement documents, nor raw pricing/commercials sheets.
_EXCLUDE_NAME_RE = re.compile(r'commercial|cost estimate', re.IGNORECASE)


def _is_excluded_source(file_path: Path) -> bool:
    if any(part.lower() == 'rfq' for part in file_path.parts):   # client RFQ subfolder
        return True
    if _EXCLUDE_NAME_RE.search(file_path.stem):                  # pricing / commercials
        return True
    return False


class ProposalLoader:

    def __init__(self):
        self.backend = PROPOSAL_VECTOR_BACKEND
        self.embeddings = get_embeddings()
        self.faiss_index = None
        self.metadata = []
        # group key -> ordered list of metadata indices, for sibling rollup
        self._parent_map: Dict[str, List[int]] = {}
        self._section_map: Dict[Tuple, int] = {}
        self._id_to_pos: Dict[str, int] = {}
        self._qdrant = None
        if self.backend == "qdrant":
            from qdrant_client import QdrantClient
            self._qdrant = QdrantClient(url=QDRANT_URL, api_key=QDRANT_API_KEY, timeout=30)

    def chunk_file(self, file_path: Path, client_name: str) -> List[Tuple[str, Dict]]:
        """Heading-aware chunking for a single source file. Returns (embed_text, meta)."""
        ext = file_path.suffix.lower()
        extractor = _EXTRACTORS.get(ext)
        if extractor is None:
            return []
        elements = extractor(file_path)
        if not elements:
            return []
        _validate_headings(elements)
        sections = _build_sections(elements)
        chunk_dicts = [c for c in _sections_to_chunks(sections)
                       if len(c["content"]) >= MIN_CHUNK_CHARS]
        if not chunk_dicts:
            return []

        sample = ' '.join(c["content"] for c in chunk_dicts[:3])
        industry = _detect_industry(sample, client_name)

        out: List[Tuple[str, Dict]] = []
        for idx, c in enumerate(chunk_dicts):
            breadcrumb = c["breadcrumb"] or c["title"]
            embed_text = (
                f"Client: {client_name} | Industry: {industry} | "
                f"Section: {breadcrumb} | File: {file_path.name}:\n\n{c['content']}"
            )
            # group key scopes the parent to THIS file (paths repeat across files)
            group_key = f"{file_path.name}::{c['parent']}" if c["parent"] else ""
            meta = {
                "client_name": client_name,
                "filename": file_path.name,
                "file_type": ext.lstrip('.'),
                "source_file": str(file_path),
                "industry": industry,
                "page_ref": c["page_ref"],
                "section_path": c["path"],
                "section_title": c["title"],
                "breadcrumb": breadcrumb,
                "parent_id": c["parent"],
                "group_key": group_key,
                "chunk_index": idx,
                "word_count": len(c["content"].split()),
                "content": c["content"],
            }
            out.append((embed_text, meta))
        return out

    def load_raw_documents(self) -> List[Tuple[str, Dict]]:
        """Read every supported document under data/raw_proposals/<client>/ (PDF,
        PPT/PPTX, DOC/DOCX — never xlsx) and chunk it heading-aware."""
        all_chunks: List[Tuple[str, Dict]] = []
        client_dirs = [d for d in RAW_PROPOSALS_DIR.iterdir() if d.is_dir()]
        if not client_dirs:
            logger.warning(f"No client folders found in {RAW_PROPOSALS_DIR}")
            return []

        for client_dir in sorted(client_dirs):
            client_name = client_dir.name
            for file_path in sorted(client_dir.rglob('*')):
                if file_path.suffix.lower() not in _EXTRACTORS:
                    continue
                if _is_excluded_source(file_path):
                    logger.info(f"Skipping non-proposal source: {file_path.name}")
                    continue
                try:
                    file_chunks = self.chunk_file(file_path, client_name)
                except Exception as e:
                    logger.error(f"Failed processing {file_path.name}: {e}")
                    continue
                if not file_chunks:
                    logger.info(f"No usable chunks from {file_path.name}")
                    continue
                all_chunks.extend(file_chunks)
                logger.info(f"  {file_path.name} -> {len(file_chunks)} chunks")

        logger.info(f"Total proposal chunks from raw files: {len(all_chunks)}")
        return all_chunks

    def create_embeddings(self, chunks: List[Tuple[str, Dict]]):
        text_only = [c[0] for c in chunks]
        logger.info(f"Embedding {len(text_only)} proposal chunks...")
        embeddings_list = self.embeddings.embed_documents(text_only)
        embeddings_array = np.array(embeddings_list).astype('float32')
        faiss.normalize_L2(embeddings_array)
        dimension = embeddings_array.shape[1]
        index = faiss.IndexFlatIP(dimension)
        index.add(embeddings_array)
        self.metadata = [c[1] for c in chunks]
        self.faiss_index = index
        self._build_parent_map()
        logger.info(f"Proposal FAISS index: {len(embeddings_list)} vectors, dim={dimension}")
        return index

    def _build_parent_map(self):
        """Index group_key -> [metadata positions] so a hit can fetch its siblings,
        and (filename, section_path) -> position so it can also fetch the parent's
        own intro chunk."""
        self._parent_map = {}
        self._section_map = {}
        for i, m in enumerate(self.metadata):
            gk = m.get("group_key")
            if gk:
                self._parent_map.setdefault(gk, []).append(i)
            self._section_map[(m.get("filename"), m.get("section_path"))] = i

    def save_index(self):
        PROPOSAL_FAISS_INDEX_PATH.parent.mkdir(parents=True, exist_ok=True)
        faiss.write_index(self.faiss_index, str(PROPOSAL_FAISS_INDEX_PATH))
        with open(PROPOSAL_METADATA_PATH, 'w', encoding='utf-8') as f:
            json.dump(self.metadata, f, indent=2, default=str)
        logger.info("Proposal index saved")

    # ── Qdrant backend ──────────────────────────────────────────────────────────
    def _ensure_qdrant_collection(self, dim: int):
        from qdrant_client.models import Distance, VectorParams
        if not self._qdrant.collection_exists(QDRANT_PROPOSAL_COLLECTION):
            self._qdrant.create_collection(
                QDRANT_PROPOSAL_COLLECTION,
                vectors_config=VectorParams(size=dim, distance=Distance.COSINE),
            )
            logger.info(f"Created Qdrant collection '{QDRANT_PROPOSAL_COLLECTION}' (dim={dim})")

    def _upsert_qdrant(self, chunks: List[Tuple[str, Dict]]) -> bool:
        from qdrant_client.models import PointStruct
        text_only = [c[0] for c in chunks]
        logger.info(f"Embedding {len(text_only)} proposal chunks for Qdrant...")
        vectors = self.embeddings.embed_documents(text_only)
        dim = len(vectors[0])
        if self._qdrant.collection_exists(QDRANT_PROPOSAL_COLLECTION):
            self._qdrant.delete_collection(QDRANT_PROPOSAL_COLLECTION)   # full snapshot
        self._ensure_qdrant_collection(dim)
        points = []
        for vec, (_, m) in zip(vectors, chunks):
            pid = str(uuid.uuid5(_QDRANT_NS, f"{m.get('source_file')}:{m.get('chunk_index')}"))
            points.append(PointStruct(id=pid, vector=[float(x) for x in vec], payload=m))
        # Upsert in batches to stay well under request-size limits.
        for i in range(0, len(points), 128):
            self._qdrant.upsert(QDRANT_PROPOSAL_COLLECTION, points=points[i:i+128])
        self.metadata = [m for _, m in chunks]
        logger.info(f"Upserted {len(points)} proposal chunks to Qdrant '{QDRANT_PROPOSAL_COLLECTION}'")
        return True

    def _load_qdrant_payloads(self):
        """Scroll all chunk payloads (metadata only — NOT vectors) into memory so the
        sibling-rollup grouped search works exactly like the FAISS path."""
        self.metadata = []
        self._id_to_pos = {}
        try:
            if not self._qdrant.collection_exists(QDRANT_PROPOSAL_COLLECTION):
                logger.warning(f"Qdrant collection '{QDRANT_PROPOSAL_COLLECTION}' does not exist yet")
                self._build_parent_map()
                return
            offset = None
            while True:
                pts, offset = self._qdrant.scroll(
                    QDRANT_PROPOSAL_COLLECTION, with_payload=True, with_vectors=False,
                    limit=256, offset=offset,
                )
                for p in pts:
                    self._id_to_pos[p.id] = len(self.metadata)
                    self.metadata.append(p.payload)
                if offset is None:
                    break
        except Exception as e:
            logger.error(f"Failed to scroll Qdrant proposal payloads: {e}")
        self._build_parent_map()
        logger.info(f"Loaded {len(self.metadata)} proposal payloads from Qdrant")

    def _search_qdrant(self, query_embedding, k: int) -> List[Dict]:
        try:
            hits = self._qdrant.query_points(
                collection_name=QDRANT_PROPOSAL_COLLECTION,
                query=[float(x) for x in query_embedding],
                limit=k, with_payload=True,
            ).points
        except Exception as e:
            logger.error(f"Qdrant proposal search failed: {e}")
            return []
        results = []
        for h in hits:
            pos = self._id_to_pos.get(h.id)
            meta = self.metadata[pos] if pos is not None else h.payload
            results.append({"similarity_score": float(h.score), "metadata": meta})
        return results

    def is_ready(self) -> bool:
        if self.backend == "qdrant":
            return bool(self.metadata)   # populated by load_index() scroll
        return bool(self.faiss_index is not None and self.metadata)

    def load_index(self):
        if self.backend == "qdrant":
            self._load_qdrant_payloads()
            return
        if PROPOSAL_FAISS_INDEX_PATH.exists():
            self.faiss_index = faiss.read_index(str(PROPOSAL_FAISS_INDEX_PATH))
        if PROPOSAL_METADATA_PATH.exists():
            with open(PROPOSAL_METADATA_PATH, 'r', encoding='utf-8') as f:
                self.metadata = json.load(f)
        self._build_parent_map()
        logger.info(f"Proposal index loaded: {len(self.metadata)} chunks")

    def search_vec(self, query_embedding, k: int = 3) -> List[Dict]:
        """Search using a pre-computed query vector (lets the caller embed once)."""
        if self.backend == "qdrant":
            return self._search_qdrant(query_embedding, k)
        if self.faiss_index is None or not self.metadata:
            return []
        qe = np.array(query_embedding).astype('float32').reshape(1, -1)
        faiss.normalize_L2(qe)
        scores, indices = self.faiss_index.search(qe, k)
        results = []
        for score, idx in zip(scores[0], indices[0]):
            if 0 <= idx < len(self.metadata):
                results.append({
                    "similarity_score": float(score),
                    "metadata": self.metadata[idx]
                })
        return results

    def search_vec_grouped(self, query_embedding, k: int = 5,
                           max_parents: int = 2, max_siblings: int = 8) -> List[Dict]:
        """Top-k search, then roll up siblings: for the strongest hits that belong to
        a parent section, pull in the other children of that parent (a metadata
        lookup, NOT bounded by k) so the LLM gets the whole section. Original hits
        keep their score; siblings are tagged is_sibling=True and inherit the
        triggering hit's score for ordering."""
        hits = self.search_vec(query_embedding, k=k)
        if not hits:
            return hits
        seen = {id(h["metadata"]): h for h in hits}  # de-dupe by metadata identity
        ordered = list(hits)
        parents_done = set()
        for h in hits:
            if len(parents_done) >= max_parents:
                break
            gk = h["metadata"].get("group_key")
            if not gk or gk in parents_done:
                continue
            sib_positions = self._parent_map.get(gk, [])
            if len(sib_positions) <= 1:
                continue
            parents_done.add(gk)
            # the parent's own intro chunk (section_path == this parent_id), if any
            extra = []
            pidx = self._section_map.get((h["metadata"].get("filename"),
                                          h["metadata"].get("parent_id")))
            if pidx is not None:
                extra.append(pidx)
            for pos in extra + sib_positions[:max_siblings]:
                meta = self.metadata[pos]
                if id(meta) in seen:
                    continue
                sib = {"similarity_score": h["similarity_score"],
                       "is_sibling": True, "metadata": meta}
                seen[id(meta)] = sib
                ordered.append(sib)
        return ordered

    def search(self, query_text: str, k: int = 3) -> List[Dict]:
        if not self.is_ready():
            return []
        return self.search_vec(self.embeddings.embed_query(query_text), k)

    def build_and_save(self) -> bool:
        chunks = self.load_raw_documents()
        if not chunks:
            logger.warning("No chunks extracted from raw proposal files")
            return False
        if self.backend == "qdrant":
            return self._upsert_qdrant(chunks)
        self.create_embeddings(chunks)
        self.save_index()
        return True
